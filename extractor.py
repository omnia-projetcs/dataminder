import os
import fitz  # PyMuPDF
import zipfile
from PIL import Image
import subprocess
import tempfile

from document_ir import DocumentIR, blocks_from_text, create_document
from logger import log_error as _log_error
from ocr_engines import ocr_image, ocr_pdf, structured_parse
from transcriber import (
    transcribe_audio, transcribe_video,
    AUDIO_EXTENSIONS, VIDEO_EXTENSIONS
)

# Module-level OCR configuration (set by main.py based on CLI args)
_ocr_engine = "tesseract"
_ocr_device = "cpu"
_ocr_lang = "en"
_ocr_dpi = 200
_ocr_max_pages = 0

# Module-level Whisper configuration (set by main.py based on CLI args)
_whisper_model = "base"
_whisper_device = "cpu"
_whisper_lang = None


def configure_ocr(engine="tesseract", device="cpu", lang="en", dpi=200, max_pages=0):
    """
    Configure the OCR engine for all extraction calls.

    Args:
        engine: "tesseract" (default, legacy) or "paddleocr" (deep learning).
        device: "cpu" or "gpu" (PaddleOCR only).
        lang: Language hint (e.g., "en", "fr", "ch").
        dpi: Resolution for rendering scanned PDF pages (default 200).
             Lower values use less RAM. Range: 72-600.
        max_pages: Maximum pages to OCR per PDF (0 = unlimited).
    """
    global _ocr_engine, _ocr_device, _ocr_lang, _ocr_dpi, _ocr_max_pages
    _ocr_engine = engine
    _ocr_device = device
    _ocr_lang = lang
    _ocr_dpi = max(72, min(600, dpi))
    _ocr_max_pages = max(0, max_pages)

    engine_name = "PaddleOCR PP-OCRv5" if engine == "paddleocr" else "Tesseract"
    extras = []
    if dpi != 200:
        extras.append(f"dpi={_ocr_dpi}")
    if max_pages > 0:
        extras.append(f"max_pages={_ocr_max_pages}")
    extra_str = f", {', '.join(extras)}" if extras else ""
    print(f"[OCR] Configured: engine={engine_name}, device={device}, lang={lang}{extra_str}")


def configure_whisper(model="base", device="cpu", lang=None):
    """
    Configure the Whisper transcription engine.

    Args:
        model: Whisper model size (tiny, base, small, medium, large-v3).
        device: "cpu" or "cuda".
        lang: Language code (e.g., "en", "fr"). None for auto-detection.
    """
    global _whisper_model, _whisper_device, _whisper_lang
    _whisper_model = model
    _whisper_device = device
    _whisper_lang = lang
    lang_display = lang if lang else "auto-detect"
    print(f"[Whisper] Configured: model={model}, device={device}, lang={lang_display}")


def extraction_config():
    """Return the active extraction configuration for cache invalidation."""
    return {
        "ocr_engine": _ocr_engine,
        "ocr_device": _ocr_device,
        "ocr_lang": _ocr_lang,
        "ocr_dpi": _ocr_dpi,
        "ocr_max_pages": _ocr_max_pages,
        "whisper_model": _whisper_model,
        "whisper_device": _whisper_device,
        "whisper_lang": _whisper_lang,
    }

def log_error(filepath, error_msg):
    _log_error(filepath, error_msg, category="EXTRACTION")

def extract_text_from_txt(path):
    with open(path, 'r', encoding='utf-8') as f:
        return f.read()

def extract_text_from_docx(path):
    from docx import Document
    doc = Document(path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pdf(path, structured=False):
    """
    Extract text from a PDF file.

    Strategy:
      1. Try direct text extraction with PyMuPDF (for native/digital PDFs)
      2. If text is sparse (scanned PDF), fall back to OCR
      3. If structured=True, use PP-StructureV3 for layout-aware extraction

    Args:
        path: Path to the PDF file.
        structured: If True, use PP-StructureV3 for structured Markdown output.
    """
    # Structured mode: use PP-StructureV3 if available
    if structured:
        result = structured_parse(path, device=_ocr_device)
        if result:
            return result
        print(f"[{path}] PP-StructureV3 unavailable or failed. Falling back to standard extraction.")

    text = ""
    # Try normal PDF extraction first
    try:
        doc = fitz.open(path)
        for page in doc:
            text += page.get_text()
        
        # If the extracted text is very short compared to the number of pages, it's likely a scanned PDF
        if len(text.strip()) < 50 * len(doc):
            print(f"[{path}] PDF appears to be scanned or contains little text. Switching to OCR...")
            text = extract_text_from_pdf_ocr(path)
    except Exception as e:
        print(f"[{path}] Error reading PDF directly: {e}. Trying OCR...")
        text = extract_text_from_pdf_ocr(path)
        
    return text

def extract_text_from_pdf_ocr(path):
    """
    OCR a scanned PDF using the configured OCR engine.

    With PaddleOCR (default):
      - Processes PDF directly (no intermediate image conversion needed)
      - Deep learning text detection + recognition (PP-OCRv5)
      - Automatic orientation correction
      - Automatic distortion/warping correction

    With Tesseract (legacy fallback):
      - Converts pages to images via pdf2image/poppler
      - Traditional LSTM-based OCR
    """
    try:
        return ocr_pdf(path, engine=_ocr_engine, device=_ocr_device, lang=_ocr_lang, dpi=_ocr_dpi, max_pages=_ocr_max_pages)
    except Exception as e:
        log_error(path, f"Failed to extract text via OCR ({_ocr_engine}). Error: {e}")
        return ""

def extract_text_from_cbz_cbr(path):
    """
    Extract text from comic book archives (CBZ/CBR) using OCR.

    With PaddleOCR: deep learning detection handles speech bubbles,
    captions, and onomatopoeia much better than Tesseract.
    """
    ext = os.path.splitext(path)[1].lower()
    text = ""
    try:
        if ext == '.cbz':
            with zipfile.ZipFile(path, 'r') as archive:
                image_files = [f for f in archive.namelist() if f.lower().endswith(('.png', '.jpg', '.jpeg', '.webp'))]
                image_files.sort()
                for img_name in image_files:
                    with archive.open(img_name) as file:
                        img = Image.open(file)
                        text += ocr_image(img, engine=_ocr_engine, device=_ocr_device, lang=_ocr_lang) + "\n"
        elif ext == '.cbr':
            import rarfile
            with rarfile.RarFile(path, 'r') as archive:
                image_files = [f for f in archive.namelist() if f.lower().endswith(('.png', '.jpg', '.jpeg', '.webp'))]
                image_files.sort()
                for img_name in image_files:
                    with archive.open(img_name) as file:
                        img = Image.open(file)
                        text += ocr_image(img, engine=_ocr_engine, device=_ocr_device, lang=_ocr_lang) + "\n"
        return text
    except Exception as e:
        log_error(path, f"Error reading archive (CBZ/CBR): {e}. Make sure 'unrar' is installed.")
        return ""

def extract_text_from_pptx(path):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        log_error(path, f"Error reading PPTX: {e}")
        return ""

def extract_text_from_doc(path):
    try:
        # Requires antiword installed on the system
        result = subprocess.run(['antiword', path], capture_output=True, text=True, check=True)
        return result.stdout
    except Exception as e:
        log_error(path, f"Error reading DOC: {e}. Make sure 'antiword' is installed.")
        return ""

def extract_text_from_excel(path):
    try:
        import pandas as pd
        text = ""
        # Read all sheets
        dfs = pd.read_excel(path, sheet_name=None)
        for sheet_name, df in dfs.items():
            text += f"--- Sheet: {sheet_name} ---\n"
            # to_string will convert the dataframe to a readable text table
            text += df.to_string(index=False) + "\n\n"
        return text
    except Exception as e:
        log_error(path, f"Error reading Excel file: {e}")
        return ""

def extract_text_from_html(path):
    try:
        from bs4 import BeautifulSoup
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f, 'html.parser')
            return soup.get_text(separator='\n', strip=True)
    except Exception as e:
        log_error(path, f"Error reading HTML: {e}")
        return ""

def extract_text_from_chm(path):
    try:
        from bs4 import BeautifulSoup
        text = ""
        # Create a temporary directory
        with tempfile.TemporaryDirectory() as temp_dir:
            # Requires libchm-bin installed on the system
            subprocess.run(['extract_chmLib', path, temp_dir], capture_output=True, check=True)
            
            # Find all HTML files in the extracted directory
            for root, _, files in os.walk(temp_dir):
                for file in files:
                    if file.lower().endswith(('.html', '.htm')):
                        file_path = os.path.join(root, file)
                        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                            soup = BeautifulSoup(f, 'html.parser')
                            extracted = soup.get_text(separator='\n', strip=True)
                            if extracted:
                                text += f"--- Section: {file} ---\n" + extracted + "\n\n"
        return text
    except Exception as e:
        log_error(path, f"Error reading CHM: {e}. Make sure 'libchm-bin' is installed.")
        return ""

def extract_text_from_epub(path):
    try:
        import ebooklib
        from bs4 import BeautifulSoup
        from ebooklib import epub
        book = epub.read_epub(path, options={'ignore_ncx': True})
        text = ""
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), 'html.parser')
            extracted = soup.get_text(separator='\n', strip=True)
            if extracted:
                text += extracted + "\n\n"
        return text
    except Exception as e:
        log_error(path, f"Error reading EPUB: {e}")
        return ""

def extract_text_from_image(path):
    """
    Extract text from a standalone image file using OCR.

    Supports: PNG, JPG, JPEG, WEBP, BMP, TIFF.
    """
    try:
        img = Image.open(path)
        return ocr_image(img, engine=_ocr_engine, device=_ocr_device, lang=_ocr_lang)
    except Exception as e:
        log_error(path, f"Error reading image for OCR: {e}")
        return ""

def extract_text_from_audio(path):
    """
    Transcribe an audio file to text using faster-whisper.

    Supports: MP3, WAV, OGG, FLAC, M4A, WMA, AAC.
    """
    try:
        return transcribe_audio(
            path,
            model_size=_whisper_model,
            device=_whisper_device,
            language=_whisper_lang
        )
    except Exception as e:
        log_error(path, f"Error transcribing audio: {e}")
        return ""

def extract_text_from_video(path):
    """
    Transcribe a video file by extracting its audio track first.

    Supports: MP4, MKV, AVI, MOV, WEBM, FLV, WMV.
    Requires ffmpeg installed on the system.
    """
    try:
        return transcribe_video(
            path,
            model_size=_whisper_model,
            device=_whisper_device,
            language=_whisper_lang
        )
    except Exception as e:
        log_error(path, f"Error transcribing video: {e}")
        return ""

def extract_text(path, structured=False):
    """
    Extract text from a document file.

    Args:
        path: Path to the document file.
        structured: If True and the file is a PDF, use PP-StructureV3
                    for layout-aware structured Markdown output.
    """
    ext = os.path.splitext(path)[1].lower()
    
    if ext in ['.txt', '.md', '.rst']:
        return extract_text_from_txt(path)
    elif ext == '.docx':
        return extract_text_from_docx(path)
    elif ext == '.pdf':
        return extract_text_from_pdf(path, structured=structured)
    elif ext in ['.cbz', '.cbr']:
        return extract_text_from_cbz_cbr(path)
    elif ext == '.pptx':
        return extract_text_from_pptx(path)
    elif ext == '.doc':
        return extract_text_from_doc(path)
    elif ext in ['.xls', '.xlsx']:
        return extract_text_from_excel(path)
    elif ext in ['.html', '.htm']:
        return extract_text_from_html(path)
    elif ext == '.chm':
        return extract_text_from_chm(path)
    elif ext == '.epub':
        return extract_text_from_epub(path)
    elif ext in ['.png', '.jpg', '.jpeg', '.webp', '.bmp', '.tiff', '.tif']:
        return extract_text_from_image(path)
    elif ext in AUDIO_EXTENSIONS:
        return extract_text_from_audio(path)
    elif ext in VIDEO_EXTENSIONS:
        return extract_text_from_video(path)
    else:
        log_error(path, f"Unsupported file extension: {ext}")
        return ""


def _pdf_document(path, source_sha256, structured=False):
    """Extract a PDF page by page while preserving provenance.

    Sparse pages are OCRed individually.  This avoids the previous all-or-nothing
    fallback where one scanned page inside an otherwise digital PDF was lost.
    """
    document = create_document(path, source_sha256=source_sha256)

    if structured:
        structured_text = structured_parse(path, device=_ocr_device)
        if structured_text:
            document.blocks = blocks_from_text(
                structured_text,
                document.id,
                extraction_method="paddle-structure-v3",
            )
            document.metadata.update(
                {"structured": True, "extraction_method": "paddle-structure-v3"}
            )
            return document
        document.diagnostics.append(
            {
                "level": "warning",
                "code": "structured_parse_failed",
                "message": "PP-StructureV3 failed; standard page extraction was used.",
            }
        )

    ocr_pages = 0
    methods = set()
    try:
        with fitz.open(path) as pdf:
            document.metadata["page_count"] = len(pdf)
            for page_index, page in enumerate(pdf):
                page_number = page_index + 1
                native_text = page.get_text()
                page_text = native_text
                method = "pymupdf"
                sparse = len(native_text.strip()) < 50
                within_ocr_limit = _ocr_max_pages == 0 or ocr_pages < _ocr_max_pages

                if sparse and within_ocr_limit:
                    try:
                        zoom = _ocr_dpi / 72.0
                        pixmap = page.get_pixmap(
                            matrix=fitz.Matrix(zoom, zoom), alpha=False
                        )
                        image = Image.frombytes(
                            "RGB",
                            (pixmap.width, pixmap.height),
                            pixmap.samples,
                        )
                        ocr_text = ocr_image(
                            image,
                            engine=_ocr_engine,
                            device=_ocr_device,
                            lang=_ocr_lang,
                        )
                        ocr_pages += 1
                        if ocr_text and ocr_text.strip():
                            page_text = ocr_text
                            method = f"{_ocr_engine}-ocr"
                        elif not native_text.strip():
                            document.diagnostics.append(
                                {
                                    "level": "warning",
                                    "code": "empty_page",
                                    "page": page_number,
                                    "message": "Native extraction and OCR returned no text.",
                                }
                            )
                    except Exception as exc:
                        document.diagnostics.append(
                            {
                                "level": "warning",
                                "code": "page_ocr_failed",
                                "page": page_number,
                                "message": str(exc),
                            }
                        )

                if page_text and page_text.strip():
                    methods.add(method)
                    document.blocks.extend(
                        blocks_from_text(
                            page_text,
                            document.id,
                            extraction_method=method,
                            page=page_number,
                            id_prefix=f"{document.id}/page/{page_number}",
                        )
                    )
    except Exception as exc:
        document.diagnostics.append(
            {
                "level": "error",
                "code": "pdf_open_failed",
                "message": str(exc),
            }
        )
        # Retain the legacy whole-document fallback for damaged PDFs.
        fallback_text = extract_text_from_pdf_ocr(path)
        document.blocks = blocks_from_text(
            fallback_text,
            document.id,
            extraction_method=f"{_ocr_engine}-ocr",
        )
        if fallback_text.strip():
            methods.add(f"{_ocr_engine}-ocr")

    document.metadata.update(
        {
            "structured": False,
            "ocr_pages": ocr_pages,
            "extraction_methods": sorted(methods),
        }
    )
    return document


def extract_document(path, structured=False, source_sha256=None) -> DocumentIR:
    """Extract *path* into the structured, engine-neutral document model."""
    document = create_document(path, source_sha256=source_sha256)
    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        return _pdf_document(path, document.source_sha256, structured=structured)

    text = extract_text(path, structured=structured)
    method_by_extension = {
        ".txt": "plain-text",
        ".md": "markdown",
        ".rst": "plain-text",
        ".docx": "python-docx",
        ".doc": "antiword",
        ".pptx": "python-pptx",
        ".xls": "pandas",
        ".xlsx": "pandas",
        ".html": "beautifulsoup",
        ".htm": "beautifulsoup",
        ".chm": "chmlib",
        ".epub": "ebooklib",
    }
    if ext in AUDIO_EXTENSIONS or ext in VIDEO_EXTENSIONS:
        method = "faster-whisper"
    elif ext in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif", ".cbz", ".cbr"}:
        method = f"{_ocr_engine}-ocr"
    else:
        method = method_by_extension.get(ext, "unknown")

    document.blocks = blocks_from_text(
        text,
        document.id,
        extraction_method=method,
    )
    document.metadata["extraction_methods"] = [method] if document.blocks else []
    return document
